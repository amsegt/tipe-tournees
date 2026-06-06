# -*- coding: utf-8 -*-
"""
make_results_pptx.py — Diaporama des RESULTATS (point b, version PowerPoint).

Cree Resultats_TIPE_13530.pptx : presentation prete a montrer, qui reprend et
commente les resultats (43 points + etalonnage Held-Karp) avec les figures.
Lit results.csv et calibration_results.csv. N'ecrase pas ta presentation existante.
"""
import csv
import os
import statistics
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

TEAL = RGBColor(0x02, 0x80, 0x90)
GREEN = RGBColor(0x2C, 0x5F, 0x2D)
ACCENT = RGBColor(0xB8, 0x50, 0x42)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)        # 16:9


def fmt(sec):
    sec = float(sec)
    m, s = divmod(int(round(sec)), 60)
    return f"{m} min {s:02d} s"


def read_csv(path):
    with open(path, newline="") as f:
        return list(csv.DictReader(f))


def add_band(slide, prs):
    """Bandeau teal en haut."""
    box = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    box.fill.solid(); box.fill.fore_color.rgb = TEAL
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def title_on_band(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = WHITE
    return tb


def bullets(slide, items, left=Inches(0.7), top=Inches(1.5),
            width=Inches(12.0), height=Inches(5.4), size=20):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for k, it in enumerate(items):
        if isinstance(it, tuple):
            text, lvl, color, bold = it
        else:
            text, lvl, color, bold = it, 0, INK, False
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.text = ("•  " if lvl == 0 else "–  ") + text
        p.level = lvl
        p.font.size = Pt(size if lvl == 0 else size - 3)
        p.font.color.rgb = color; p.font.bold = bold
        p.space_after = Pt(8)
    return tb


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_image_fit(slide, path, left, top, max_w, max_h):
    """Insere une image en respectant le ratio, dans la boite (max_w, max_h)."""
    from PIL import Image
    try:
        with Image.open(path) as im:
            iw, ih = im.size
    except Exception:
        slide.shapes.add_picture(path, left, top, width=max_w)
        return
    ratio = min(max_w / iw, max_h / ih)
    w = int(iw * ratio); h = int(ih * ratio)
    off_l = left + int((max_w - w) / 2)
    slide.shapes.add_picture(path, off_l, top, width=Emu(w), height=Emu(h))


def caption(slide, text, top=Inches(6.9)):
    tb = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.5))
    p = tb.text_frame.paragraphs[0]; p.text = text
    p.font.size = Pt(13); p.font.color.rgb = GREY; p.font.italic = True


# ---- tableaux ----
def results_table(slide, results, best):
    rows = len(results) + 1
    tbl_shape = slide.shapes.add_table(rows, 3, Inches(0.9), Inches(1.7),
                                       Inches(11.5), Inches(0.5 * rows))
    table = tbl_shape.table
    table.columns[0].width = Inches(6.7)
    table.columns[1].width = Inches(2.6)
    table.columns[2].width = Inches(2.2)
    heads = ["Methode", "Temps de cycle", "Calcul (s)"]
    for c, h in enumerate(heads):
        cell = table.cell(0, c); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True; para.font.color.rgb = WHITE; para.font.size = Pt(16)
    for r, row in enumerate(results, start=1):
        is_best = row is best
        vals = [row["methode"], fmt(row["longueur_secondes"]),
                f'{float(row["temps_calcul_s"]):.4f}']
        for c, v in enumerate(vals):
            cell = table.cell(r, c); cell.text = v
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(15)
            para.font.color.rgb = ACCENT if is_best else INK
            para.font.bold = bool(is_best)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFB, 0xEF, 0xEC) if is_best else WHITE


def calib_table(slide, g0, gb, g2):
    data = [("Plus proche voisin (depot)", g0, False),
            ("Plus proche voisin (meilleur depart)", gb, False),
            ("2-opt (sur meilleur PPV)", g2, True)]
    tbl_shape = slide.shapes.add_table(4, 4, Inches(0.9), Inches(1.7),
                                       Inches(11.5), Inches(2.2))
    table = tbl_shape.table
    table.columns[0].width = Inches(5.9)
    for i, w in enumerate([2.0, 1.9, 1.7], start=1):
        table.columns[i].width = Inches(w)
    for c, h in enumerate(["Heuristique", "ecart moyen", "min", "max"]):
        cell = table.cell(0, c); cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = TEAL
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True; para.font.color.rgb = WHITE; para.font.size = Pt(16)
    for r, (label, g, hl) in enumerate(data, start=1):
        vals = [label, f"+{statistics.mean(g):.1f} %",
                f"+{min(g):.1f} %", f"+{max(g):.1f} %"]
        for c, v in enumerate(vals):
            cell = table.cell(r, c); cell.text = v
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(15)
            para.font.color.rgb = ACCENT if hl else INK
            para.font.bold = bool(hl)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFB, 0xEF, 0xEC) if hl else WHITE


def main():
    results = read_csv("results.csv")
    calib = read_csv("calibration_results.csv")
    best = min(results, key=lambda r: float(r["longueur_secondes"]))
    g0 = [float(r["ecart_ppv_depot_%"]) for r in calib]
    gb = [float(r["ecart_ppv_best_%"]) for r in calib]
    g2 = [float(r["ecart_2opt_%"]) for r in calib]

    prs = Presentation()
    prs.slide_width = SW; prs.slide_height = SH

    # --- 1. Titre ---
    s = blank(prs)
    bg = s.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid(); bg.fill.fore_color.rgb = TEAL; bg.line.fill.background()
    bg.shadow.inherit = False
    tb = s.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.3), Inches(2.5))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = "Tournees de collecte : comparaison des algorithmes"
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = "TIPE 13530 — Brahim Amsegt — Plus proche voisin / 2-opt / Held-Karp"
    p2.font.size = Pt(20); p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER
    p3 = tf.add_paragraph()
    p3.text = "Resultats : 2-opt quasi-optimal (+0,9 % a l'optimum exact)"
    p3.font.size = Pt(18); p3.font.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    p3.font.italic = True; p3.alignment = PP_ALIGN.CENTER

    # --- 2. Probleme & modelisation ---
    s = blank(prs); add_band(s, prs); title_on_band(s, "1. Le probleme : un TSP")
    bullets(s, [
        ("Collecte = tournee qui visite tous les points et revient au depot.", 0, INK, True),
        ("Modele : graphe complet pondere G = (V, E, w).", 0, INK, False),
        ("V = points de collecte ; w(u, v) = temps de trajet estime.", 1, GREY, False),
        ("On cherche le CYCLE HAMILTONIEN de poids minimal.", 0, INK, False),
        ("Probleme NP-difficile : aucun algorithme polynomial exact connu.", 0, ACCENT, True),
        ("(n-1)!/2 cycles possibles : > 43 milliards des 15 points.", 1, GREY, False),
        ("Trois algorithmes, trois compromis vitesse / qualite / exactitude.", 0, INK, True),
    ])

    # --- 3. Les trois algorithmes ---
    s = blank(prs); add_band(s, prs); title_on_band(s, "2. Les trois algorithmes")
    bullets(s, [
        ("Plus proche voisin (PPV) — glouton, O(n^2).", 0, GREEN, True),
        ("A chaque etape, le point non visite le plus proche. Rapide, myope.", 1, GREY, False),
        ("2-opt — amelioration locale, O(n^2) par passe.", 0, ACCENT, True),
        ("Decroise le tour du PPV tant qu'un echange le raccourcit. Optimum local.", 1, GREY, False),
        ("Held-Karp — exact, programmation dynamique, O(n^2 . 2^n).", 0, TEAL, True),
        ("Optimum garanti, mais mur memoire vers n = 20 : reserve a l'etalonnage.", 1, GREY, False),
    ])

    # --- 4. Resultats 43 points ---
    s = blank(prs); add_band(s, prs)
    title_on_band(s, "3. Resultats sur l'instance reelle (43 points)")
    results_table(s, results, best)
    bullets(s, [
        ("2-opt = meilleur cycle : 13 min 23 s (-26,5 % vs PPV-depot).", 0, ACCENT, True),
        ("PPV depend fortement du point de depart (~10 % de gain au meilleur depart).", 0, INK, False),
        ("2-opt calcule meme plus vite que le PPV-meilleur-depart, tout en gagnant.", 0, INK, False),
    ], top=Inches(4.7), height=Inches(2.4), size=18)

    # --- 5. Etalonnage Held-Karp (tableau) ---
    s = blank(prs); add_band(s, prs)
    title_on_band(s, "4. Etalonnage contre l'optimum exact (Held-Karp)")
    calib_table(s, g0, gb, g2)
    bullets(s, [
        ("7 sous-quartiers de 14 points : Held-Karp donne l'optimum EXACT.", 0, INK, True),
        ("2-opt quasi-optimal : +0,9 % en moyenne, optimum atteint 2 fois sur 7.", 0, ACCENT, True),
        ("Le plus proche voisin peut deriver jusqu'a +24,8 %.", 0, INK, False),
    ], top=Inches(4.2), height=Inches(2.8), size=18)

    # --- 6. Figure etalonnage ---
    s = blank(prs); add_band(s, prs)
    title_on_band(s, "4. Etalonnage — distribution des ecarts")
    if os.path.exists("calibration_gaps.png"):
        add_image_fit(s, "calibration_gaps.png", Inches(1.4), Inches(1.35),
                      Inches(10.5), Inches(5.3))
    caption(s, "Ecart a l'optimum exact par sous-quartier : le 2-opt (rouge) reste colle a 0 %.")

    # --- 7. Meilleur cycle ---
    s = blank(prs); add_band(s, prs)
    title_on_band(s, "5. Le meilleur cycle trouve (2-opt)")
    if os.path.exists("best_tour.png"):
        add_image_fit(s, "best_tour.png", Inches(2.6), Inches(1.35),
                      Inches(8.1), Inches(5.3))
    caption(s, "Cycle hamiltonien retenu sur les 43 points. Carre = depot (depart/arrivee).")

    # --- 8. Graphe + matrice ---
    s = blank(prs); add_band(s, prs)
    title_on_band(s, "6. Le graphe et sa matrice d'adjacence")
    if os.path.exists("graph_view.png"):
        add_image_fit(s, "graph_view.png", Inches(0.4), Inches(1.35),
                      Inches(6.2), Inches(5.2))
    if os.path.exists("adjacency_heatmap.png"):
        add_image_fit(s, "adjacency_heatmap.png", Inches(6.8), Inches(1.35),
                      Inches(6.2), Inches(5.2))
    caption(s, "A gauche : le graphe pondere. A droite : la matrice 43x43 (graphe complet).")

    # --- 9. Conclusion / rigueur ---
    s = blank(prs); add_band(s, prs); title_on_band(s, "7. Conclusion & limites")
    bullets(s, [
        ("Classement net et robuste : 2-opt > PPV-best > PPV-depot.", 0, ACCENT, True),
        ("2-opt = bon compromis : quasi-optimal et tres rapide.", 0, INK, True),
        ("Limites assumees :", 0, INK, True),
        ("temps en hypothese forte (vol d'oiseau, vitesse constante) = borne.", 1, GREY, False),
        ("ecart mesure sur petites instances, extrapole a 43 points avec prudence.", 1, GREY, False),
        ("TSP = simplification assumee d'un CARP.", 1, GREY, False),
        ("Prolongement : matrice sur reseau routier reel (OSM) -> test de robustesse.", 0, TEAL, True),
    ])

    prs.save("Resultats_TIPE_13530.pptx")
    print("PPTX ecrit -> Resultats_TIPE_13530.pptx")


if __name__ == "__main__":
    main()
